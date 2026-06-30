"""Document text extraction from various file formats."""
import io
import os
import pypdf
import docx
import pptx
import openpyxl


class DocumentParser:

    SUPPORTED_MIMES = {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "text/plain": "txt",
        "text/markdown": "md",
    }

    EXT_MAP = {
        ".pdf": "pdf", ".docx": "docx", ".pptx": "pptx",
        ".xlsx": "xlsx", ".txt": "txt", ".md": "md",
        ".html": "html", ".htm": "html",
    }

    def parse(self, content: bytes, mime_type: str, filename: str = "") -> str:
        import logging
        log = logging.getLogger(__name__)
        log.info(f"PARSER: mime_type={mime_type!r}, filename={filename!r}, content_len={len(content)}")
        mime_type = mime_type.split(";")[0].strip()
        ext = self.SUPPORTED_MIMES.get(mime_type, "")
        if not ext and filename:
            _, fext = os.path.splitext(filename.lower())
            ext = self.EXT_MAP.get(fext, "")
        if ext == "pdf":
            return self._parse_pdf(content)
        elif ext == "docx":
            return self._parse_docx(content)
        elif ext == "pptx":
            return self._parse_pptx(content)
        elif ext == "xlsx":
            return self._parse_xlsx(content)
        else:
            return content.decode("utf-8", errors="replace")

    def _parse_pdf(self, content: bytes) -> str:
        """Parse PDF: try pypdf first, then MinerU (OCR), then UTF-8 fallback."""
        # Try pypdf (fast, for text-based PDFs)
        try:
            reader = pypdf.PdfReader(io.BytesIO(content))
            pages_text = []
            for page in reader.pages:
                try:
                    t = page.extract_text()
                    if t and t.strip():
                        pages_text.append(t.strip())
                except Exception:
                    continue
            if pages_text:
                return "\n\n".join(pages_text)
        except Exception as e:
            logging.getLogger(__name__).warning(f"pypdf failed: {e}")

        # Try MinerU (for scanned PDFs / OCR)
        try:
            return self._parse_pdf_with_mineru(content)
        except ImportError:
            pass
        except Exception as e:
            logging.getLogger(__name__).warning(f"MinerU failed: {e}, falling back to UTF-8")

        # Final fallback
        return content.decode("utf-8", errors="replace")

    def _parse_pdf_with_mineru(self, content: bytes) -> str:
        """Parse PDF using MinerU (magic-pdf) for OCR/layout analysis."""
        from magic_pdf.pipe import UNIPipe
        from magic_pdf.rw import AbsReaderWriter

        jso_useful_key = {"pdf_bytes": content}
        pipe = UNIPipe(content, jso_useful_key)
        pipe.pipe_classify()
        pipe.pipe_analyze()
        pipe.pipe_parse()
        result = pipe.get_result()
        md_content = result.get("md_content", "") or ""

        # Clean up markdown formatting, keep plain text
        text = re.sub(r"[#*_~`>|\\-]+", "", md_content)
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        return text if text else ""
    def _parse_docx(self, content: bytes) -> str:
        doc = docx.Document(io.BytesIO(content))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _parse_pptx(self, content: bytes) -> str:
        prs = pptx.Presentation(io.BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
        return "\n\n".join(texts)

    def _parse_xlsx(self, content: bytes) -> str:
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
        rows = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = "\t".join(cells)
                if line.strip():
                    rows.append(f"[{sheet.title}] {line}")
        return "\n".join(rows)
